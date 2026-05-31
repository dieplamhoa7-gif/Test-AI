from pathlib import Path
import json
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

model=json.loads(Path('mwg_segment_model_public.json').read_text(encoding='utf-8'))
shares_m=model['MWG']['shares_m']
net_debt_bn=model['MWG']['net_debt_bn']
price=model['MWG']['price']

# build scenarios
scenarios={
    'Bear': {'BHX':30075.0,'DMX':32000.0,'TGDD':37868.0,'Other':9467.0},
    'Base': {'BHX':37837.5,'DMX':45000.0,'TGDD':21258.0,'Other':5314.5},
    'Bull': {'BHX':45600.0,'DMX':60000.0,'TGDD':3048.0,'Other':762.0},
}
# note: Bear TGDD+Other and Bull TGDD+Other reflect prior implied buckets split 80/20.
# compute equity / target price
for name,v in scenarios.items():
    ev=sum(v.values())
    eq=ev-net_debt_bn
    tp=eq*1000/shares_m
    upside=tp/price-1
    v['MWG_EV']=ev; v['MWG_Equity']=eq; v['TargetPrice']=tp; v['Upside']=upside

# Excel
wb=Workbook()
ws=wb.active
ws.title='MWG_SOTP'
blue=PatternFill('solid', fgColor='1F4E78')
for c in ws[1]:
    c.font=Font(bold=True,color='FFFFFF'); c.fill=blue
headers=['Metric','Bear','Base','Bull','Notes']
for i,h in enumerate(headers,1): ws.cell(1,i).value=h
rows=[
    ['BHX value (bn)', scenarios['Bear']['BHX'], scenarios['Base']['BHX'], scenarios['Bull']['BHX'], 'BHX blend EV/Sales + P/E from public plan inputs'],
    ['DMX value (bn)', scenarios['Bear']['DMX'], scenarios['Base']['DMX'], scenarios['Bull']['DMX'], 'IPO/public narrative anchored scenario'],
    ['TGDD value (bn)', scenarios['Bear']['TGDD'], scenarios['Base']['TGDD'], scenarios['Bull']['TGDD'], 'Implied remainder split; assumption'],
    ['Other value (bn)', scenarios['Bear']['Other'], scenarios['Base']['Other'], scenarios['Bull']['Other'], 'Implied remainder split; assumption'],
    ['MWG EV (bn)', scenarios['Bear']['MWG_EV'], scenarios['Base']['MWG_EV'], scenarios['Bull']['MWG_EV'], 'Sum of segment values'],
    ['Net debt (bn)', net_debt_bn, net_debt_bn, net_debt_bn, 'Model placeholder; update with latest BS'],
    ['MWG equity value (bn)', scenarios['Bear']['MWG_Equity'], scenarios['Base']['MWG_Equity'], scenarios['Bull']['MWG_Equity'], 'EV - net debt'],
    ['Shares (m)', shares_m, shares_m, shares_m, 'Model placeholder'],
    ['Target price (VND/share)', scenarios['Bear']['TargetPrice'], scenarios['Base']['TargetPrice'], scenarios['Bull']['TargetPrice'], 'Equity value / shares'],
    ['Upside vs 70,000', scenarios['Bear']['Upside'], scenarios['Base']['Upside'], scenarios['Bull']['Upside'], 'Using current placeholder price 70,000'],
]
for r,row in enumerate(rows,2):
    for c,val in enumerate(row,1): ws.cell(r,c).value=val
for col in ['A','B','C','D','E']: ws.column_dimensions[col].width=24 if col!='E' else 58
for cell in ws['A']: cell.font=Font(bold=True)
for row in ws.iter_rows():
    for c in row: c.alignment=Alignment(wrap_text=True,vertical='top')
for r in range(2,12):
    for c in [2,3,4]:
        if r==10: ws.cell(r,c).number_format='0.0%'
        elif r==9: ws.cell(r,c).number_format='#,##0'
        else: ws.cell(r,c).number_format='#,##0.0'
ws2=wb.create_sheet('Assumption_Log')
ws2.append(['Item','Status'])
for item in [
    ('MWG 2025 KQKD','Sourced from CafeF HTML'),
    ('BHX 2026 revenue/profit plan','Public CafeF article'),
    ('BHX Q1/2026 revenue/profit','Public CafeF article'),
    ('DMX Q1/2026 profit','Public MWG/CafeF article'),
    ('DMX full-year standalone valuation','Assumption/range, not fully sourced FS'),
    ('TGDD standalone valuation','Assumption via implied remainder'),
    ('Net debt','Placeholder, needs exact BS update'),
]: ws2.append(list(item))
for col in ['A','B']: ws2.column_dimensions[col].width=42
xlsx=Path('MWG_SOTP_Scenarios.xlsx'); wb.save(xlsx)

# update PPT with target price slide
ppt=Path('MWG_BHX_DMX_SOTP_Report.pptx')
prs=Presentation(str(ppt))
slide=prs.slides.add_slide(prs.slide_layouts[6])
BLUE=RGBColor(31,78,121); DARK=RGBColor(34,34,34)
tb=slide.shapes.add_textbox(Inches(0.5),Inches(0.3),Inches(12.2),Inches(0.8))
p=tb.text_frame.paragraphs[0]; r=p.add_run(); r.text='Target price MWG theo SOTP scenario'; r.font.bold=True; r.font.size=Pt(24); r.font.color.rgb=BLUE
# table
rows_p=[['Kịch bản','Target price','Upside vs 70k'],
        ['Bear', f"{scenarios['Bear']['TargetPrice']:,.0f}".replace(',','.'), f"{scenarios['Bear']['Upside']*100:.1f}%"],
        ['Base', f"{scenarios['Base']['TargetPrice']:,.0f}".replace(',','.'), f"{scenarios['Base']['Upside']*100:.1f}%"],
        ['Bull', f"{scenarios['Bull']['TargetPrice']:,.0f}".replace(',','.'), f"{scenarios['Bull']['Upside']*100:.1f}%"],]
table=slide.shapes.add_table(len(rows_p),3,Inches(0.9),Inches(1.5),Inches(6.5),Inches(2.2)).table
for i,row in enumerate(rows_p):
    for j,val in enumerate(row):
        cell=table.cell(i,j); cell.text=str(val)
        for par in cell.text_frame.paragraphs:
            par.font.size=Pt(14); par.font.color.rgb=DARK
            if i==0: par.font.bold=True; par.font.color.rgb=BLUE
# bullets
bx=slide.shapes.add_textbox(Inches(7.8),Inches(1.4),Inches(4.6),Inches(4.4))
tf=bx.text_frame
for idx,text in enumerate([
    'Base case target price chỉ mang tính tham chiếu vì net debt và standalone segment FS chưa khóa tuyệt đối.',
    'Bear/Base/Bull được tạo để nhìn độ nhạy valuation khi BHX/ĐMX thay đổi.',
    'Nếu cập nhật net debt chính xác và standalone ĐMX/BHX sạch hơn, target price sẽ thay đổi đáng kể.'
]):
    p=tf.paragraphs[0] if idx==0 else tf.add_paragraph(); p.text=text; p.font.size=Pt(16); p.font.color.rgb=DARK
prs.save(ppt)

print(xlsx.resolve())
print(ppt.resolve())
for k,v in scenarios.items():
    print(k, round(v['TargetPrice'],0), f"{v['Upside']*100:.1f}%")
