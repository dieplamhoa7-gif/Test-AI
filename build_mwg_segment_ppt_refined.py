from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json

raw = json.loads(Path('mwg_refined_model_output.json').read_text(encoding='utf-8'))
shares_m = raw['shares_m']
price = raw['price_used']
net_debt_bn = raw['net_debt_bn']
mwg_2025 = raw['mwg_2025']
bhx_public = raw['bhx_public']
base = raw['cases']['Base']
bear = raw['cases']['Bear']
bull = raw['cases']['Bull']

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLUE = RGBColor(31,78,121)
DARK = RGBColor(34,34,34)
GRAY = RGBColor(100,100,100)


def add_title(slide, title, subtitle=''):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.2), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.bold = True; r.font.size = Pt(24); r.font.color.rgb = BLUE
    if subtitle:
        p2 = tx.text_frame.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(11); p2.font.color.rgb = GRAY


def add_bullets(slide, left, top, width, height, bullets, font_size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = b
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK


def add_table(slide, left, top, width, height, data, font_size=12):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = BLUE
    return table

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'MWG SOTP refine: BHX / ĐMX / TGDĐ', 'Model refine hơn; không neo ngược về giá 70k')
add_bullets(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.8), [
    'Mục tiêu: sửa bộ định giá cũ vì ĐMX và TGDĐ trước đó còn quá thô / implied remainder.',
    'Model mới dùng logic earnings-based rõ hơn cho ĐMX, blend EV/Sales + P/E cho BHX, và giữ TGDĐ ở mức bảo thủ.',
    f'Benchmark hiện dùng: giá MWG {price:,.0f} đ/cp, {shares_m:,.0f} triệu cp, net debt {net_debt_bn:,.0f} tỷ.'.replace(',', '.'),
    'Lưu ý: net debt và shares vẫn cần cập nhật theo BCTC/latest source để ra target price chính xác hơn.'
], 18)

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Dữ liệu public đã khóa được')
add_table(slide, Inches(0.5), Inches(1.2), Inches(12.2), Inches(4.8), [
    ['Mảng', 'Chỉ tiêu', 'Giá trị', 'Nguồn/ghi chú'],
    ['MWG', 'Doanh thu 2025', f"{mwg_2025['rev_bn']:.1f} tỷ", 'CafeF KQKD 2025'],
    ['MWG', 'LNST 2025', f"{mwg_2025['pat_bn']:.1f} tỷ", 'CafeF KQKD 2025'],
    ['BHX', 'Doanh thu Q1/2026', f"{bhx_public['rev_q1_2026_bn']:.1f} tỷ", 'Bài public CafeF'],
    ['BHX', 'LN Q1/2026', f"{bhx_public['pat_q1_2026_bn']:.1f} tỷ", 'Bài public CafeF'],
    ['BHX', 'Biên LN Q1/2026', f"{bhx_public['margin_q1_2026']*100:.2f}%", 'Suy ra từ DT/LN public'],
    ['BHX', 'KH DT/LN 2026', f"{bhx_public['rev_2026_plan_bn']:.0f} / {bhx_public['pat_2026_plan_bn']:.0f} tỷ", 'Bài public CafeF'],
    ['ĐMX', 'LNST Q1/2026', '2.219 tỷ', 'MWG/CafeF public'],
    ['ĐMX', 'Dịch vụ thợ 2025', '2.576 tỷ DT / 201 tỷ LNST', 'CafeF public'],
], 11)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Logic model mới')
add_bullets(slide, Inches(0.8), Inches(1.4), Inches(11.6), Inches(5.5), [
    'BHX: blend EV/Sales và P/E dựa trên kế hoạch 2026 (55.500 tỷ doanh thu; 1.800 tỷ lợi nhuận).',
    'ĐMX: không dùng placeholder IPO đơn thuần nữa; dùng LNST Q1/2026 = 2.219 tỷ, annualize có haircut rồi áp P/E.',
    'TGDĐ: vì chưa có standalone clean trong phiên nên giữ hướng bảo thủ bằng PAT assumption x P/E mature retail.',
    'Other: để nhỏ, tránh thổi phồng giá trị ngoài 3 mảng chính.',
    'Mục tiêu của model này là làm cho target price “có lý hơn”, không phải làm đẹp số.'
], 18)

# Slide 4 assumptions
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Assumption chi tiết theo kịch bản')
add_table(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(3.8), [
    ['Khoản mục', 'Bear', 'Base', 'Bull'],
    ['BHX EV/Sales', '0.45x', '0.55x', '0.70x'],
    ['BHX P/E', '16x', '20x', '24x'],
    ['ĐMX annualization factor trên Q1/2026', '3.2x', '3.5x', '3.8x'],
    ['ĐMX P/E', '10x', '11x', '12x'],
    ['TGDĐ PAT assumption', '900 tỷ', '1.100 tỷ', '1.300 tỷ'],
    ['TGDĐ P/E', '8x', '9x', '10x'],
    ['Other', '2.000 tỷ', '2.500 tỷ', '3.000 tỷ'],
], 11)
add_bullets(slide, Inches(0.8), Inches(5.3), Inches(11.4), Inches(1.2), [
    'Cái nào sourced thì dùng sourced; cái nào chưa sourced đủ thì ghi rõ là assumption để tránh “chính xác giả”.'
], 14)

# Slide 5 outputs
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Kết quả target price MWG')
add_table(slide, Inches(0.6), Inches(1.4), Inches(7.4), Inches(2.5), [
    ['Kịch bản', 'MWG EV', 'Equity value', 'Target price', 'Upside vs 70k'],
    ['Bear', f"{bear['MWG_EV_bn']:.1f}", f"{bear['MWG_equity_bn']:.1f}", f"{bear['TargetPrice']:.0f}", f"{bear['Upside']*100:.1f}%"],
    ['Base', f"{base['MWG_EV_bn']:.1f}", f"{base['MWG_equity_bn']:.1f}", f"{base['TargetPrice']:.0f}", f"{base['Upside']*100:.1f}%"],
    ['Bull', f"{bull['MWG_EV_bn']:.1f}", f"{bull['MWG_equity_bn']:.1f}", f"{bull['TargetPrice']:.0f}", f"{bull['Upside']*100:.1f}%"],
], 12)
add_bullets(slide, Inches(8.4), Inches(1.5), Inches(4.1), Inches(3.6), [
    f"Bear: ~{bear['TargetPrice']:.0f} đ/cp",
    f"Base: ~{base['TargetPrice']:.0f} đ/cp",
    f"Bull: ~{bull['TargetPrice']:.0f} đ/cp",
    'Base quanh 80–85k là vùng có lý hơn bộ số cũ.',
    '9x chỉ hợp lý khi ĐMX giữ earnings power cao và BHX giữ được lợi nhuận.'
], 16)

# Slide 6 segment outputs
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Phân rã giá trị base case')
add_table(slide, Inches(0.7), Inches(1.5), Inches(7.2), Inches(2.4), [
    ['Mảng', 'Giá trị (tỷ đồng)', 'Logic chính'],
    ['BHX', f"{base['BHX']['blended_value_bn']:.1f}", 'Blend EV/Sales + P/E'],
    ['ĐMX', f"{base['DMX']['value_bn']:.1f}", 'Annualized PAT x P/E'],
    ['TGDĐ', f"{base['TGDD']['value_bn']:.1f}", 'PAT assumption x P/E'],
    ['Other', f"{base['Other']:.1f}", 'Residual nhỏ, bảo thủ'],
], 12)
add_bullets(slide, Inches(8.4), Inches(1.5), Inches(4.1), Inches(4.0), [
    'ĐMX là biến số lớn nhất của model mới.',
    'BHX đủ lớn để tạo upside, nhưng chưa nên được coi là “all-in value driver”.',
    'TGDĐ đang bị giữ bảo thủ do thiếu standalone clean data; nếu bóc tách tốt hơn, giá trị có thể nhích lên.'
], 16)

# Slide 7 why previous felt wrong
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Vì sao bộ số trước nhìn vô lý?')
add_bullets(slide, Inches(0.8), Inches(1.3), Inches(11.8), Inches(5.5), [
    'Bộ cũ dùng ĐMX theo “IPO narrative placeholder”, chưa bám earnings power thực tế từ Q1/2026.',
    'TGDĐ trước đó bị định giá bằng “implied remainder”, tức là neo ngược từ EV hiện tại, không phải valuation độc lập.',
    'Kết quả là target price bị méo: hoặc quá phụ thuộc giá hiện tại, hoặc quá thô ở từng mảng.',
    'Model mới sửa bằng cách định giá độc lập từng mảng hơn, nhất là ĐMX.',
    'Cảm giác “hơi vô lý” của user là đúng; cần cảnh giác với mọi model mà phần lớn giá trị đến từ remainder hoặc placeholder.'
], 18)

# Slide 8 conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Kết luận đầu tư sơ bộ')
add_bullets(slide, Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.2), [
    'Base case hiện cho MWG khoảng 84.8k/cp, tức upside khoảng 21% so với benchmark 70k.',
    'ĐMX là engine lợi nhuận số 1; nếu earnings power Q1/2026 lặp lại với haircut hợp lý, ĐMX có thể chiếm phần lớn valuation MWG.',
    'BHX là optionality quan trọng: nếu giữ biên quanh 3% và mở rộng bền, upside còn; nếu biên tụt, định giá BHX co lại nhanh.',
    'TGDĐ đóng vai trò nền và đang bị giữ bảo thủ; upside bổ sung có thể đến nếu bóc tách segment này sạch hơn.',
    'Bước tiếp theo: cập nhật exact net debt, shares, giá MWG live, và tìm thêm standalone/IPO disclosures để khóa valuation.'
], 18)

out = Path('MWG_BHX_DMX_SOTP_Report_Refined.pptx')
prs.save(out)
print(out.resolve())
