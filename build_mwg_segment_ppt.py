from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json

# Public-input segment valuation model
shares_m = 1463.0
price = 70000.0
net_debt_bn = 7000.0
mwg_market_cap_bn = shares_m * price / 1000.0
mwg_ev_bn = mwg_market_cap_bn + net_debt_bn

# Segment assumptions / public data
bhx_rev_2026_bn = 55500.0
bhx_profit_2026_bn = 1800.0
bhx_margin = bhx_profit_2026_bn / bhx_rev_2026_bn
bhx_ev_sales = 0.65
bhx_pe = 22.0
bhx_value_bn = (bhx_rev_2026_bn * bhx_ev_sales + bhx_profit_2026_bn * bhx_pe) / 2.0

dmx_value_bn = 45000.0
# implied TGDD + other from base SOTP
other_value_bn = mwg_ev_bn - bhx_value_bn - dmx_value_bn
# split implied remainder between TGDD and other using judgment
# because TGDD is still substantial while other chains are smaller
# 80/20 split of remainder
TGDD_share = 0.80
tgdd_value_bn = other_value_bn * TGDD_share
other_chains_bn = other_value_bn * 0.20

# Save model json
model = {
    'MWG': {'market_cap_bn': mwg_market_cap_bn, 'ev_bn': mwg_ev_bn, 'shares_m': shares_m, 'price': price, 'net_debt_bn': net_debt_bn},
    'BHX': {'rev_2026_bn': bhx_rev_2026_bn, 'profit_2026_bn': bhx_profit_2026_bn, 'net_margin': bhx_margin, 'ev_sales': bhx_ev_sales, 'pe': bhx_pe, 'valuation_bn': bhx_value_bn},
    'DMX': {'valuation_bn': dmx_value_bn},
    'TGDD': {'valuation_bn': tgdd_value_bn},
    'Other': {'valuation_bn': other_chains_bn},
    'notes': [
        'BHX value blends EV/Sales and P/E using public 2026 plan.',
        'DMX base valuation uses IPO/public narrative placeholder range midpoint previously used in report.',
        'TGDD valuation is implied remainder split 80/20 from TGDD+other bucket; this is an assumption, not sourced standalone financial valuation.',
        'All values are preliminary and depend on updating price, net debt, and cleaner standalone segment disclosures.'
    ]
}
Path('mwg_segment_model_public.json').write_text(json.dumps(model, ensure_ascii=False, indent=2), encoding='utf-8')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(31,78,121)
DARK = RGBColor(34,34,34)
GRAY = RGBColor(100,100,100)
GREEN = RGBColor(34,139,34)
ORANGE = RGBColor(227,108,10)


def add_title(slide, title, subtitle=''):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.2), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.bold = True; r.font.size = Pt(24); r.font.color.rgb = BLUE
    if subtitle:
        p2 = tx.text_frame.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(11); p2.font.color.rgb = GRAY


def add_bullets(slide, left, top, width, height, bullets, font_size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = b
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK


def add_table(slide, left, top, width, height, data):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r,c)
            cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = BLUE
    return table

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Định giá SOTP MWG / ĐMX / BHX / TGDĐ', 'Public-input model; giả định được ghi rõ')
add_bullets(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.8), [
    'Mục tiêu: chạy nhanh định giá từng mảng ĐMX, BHX, TGDĐ để nhìn đóng góp vào MWG.',
    'Nguồn public đã dùng: CafeF KQKD MWG 2025; bài công khai về BHX và ĐMX IPO/operating update; dữ liệu công khai trên MWG/CafeF.',
    'Phần nào không có standalone financials sạch sẽ được đánh dấu là assumption.',
    f'MWG giả định tại giá {price:,.0f} đ/cp, {shares_m:,.0f} triệu cp, net debt {net_debt_bn:,.0f} tỷ => EV {mwg_ev_bn:,.0f} tỷ.'.replace(',', '.')
], 18)

# Slide 2: sourced data
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Dữ liệu public đã xác nhận')
add_table(slide, Inches(0.5), Inches(1.2), Inches(12.2), Inches(4.8), [
    ['Mảng','Chỉ tiêu','Giá trị','Ghi chú'],
    ['MWG','Doanh thu thuần 2025','155.928 tỷ','CafeF KQKD 2025'],
    ['MWG','LNST 2025','7.073 tỷ','CafeF KQKD 2025'],
    ['BHX','Doanh thu Q1/2026','13.100 tỷ','Bài CafeF public'],
    ['BHX','LN Q1/2026','~400 tỷ','Bài CafeF public'],
    ['BHX','KH doanh thu 2026','55.500 tỷ','Bài CafeF public'],
    ['BHX','KH lợi nhuận 2026','1.800 tỷ','Bài CafeF public'],
    ['ĐMX','LNST Q1/2026','~2.219 tỷ','Nguồn MWG/CafeF public'],
    ['ĐMX','Dịch vụ thợ 2025 DT/LNST','2.576 / 201 tỷ','Bài CafeF public'],
])

# Slide 3: assumptions
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Giả định định giá')
add_table(slide, Inches(0.6), Inches(1.3), Inches(11.8), Inches(3.0), [
    ['Khoản mục','Giả định base','Loại'],
    ['BHX EV/Sales','0,65x','Assumption'],
    ['BHX P/E','22x','Assumption'],
    ['ĐMX valuation','45.000 tỷ','Base placeholder theo IPO/public narrative'],
    ['TGDĐ valuation','Implied from remainder','Assumption'],
    ['TGDĐ / Other split','80% / 20%','Assumption'],
    ['MWG net debt','7.000 tỷ','Model placeholder - cần update BS'],
])
add_bullets(slide, Inches(0.8), Inches(4.7), Inches(11.4), Inches(1.8), [
    'BHX dùng blend EV/Sales và P/E vì đang ở pha scale + profitability transition.',
    'ĐMX chưa có bộ standalone full-year sạch trong phiên nên dùng range IPO/public narrative, lấy base 45.000 tỷ.',
    'TGDĐ chưa được tách riêng hoàn toàn bằng financials public, nên định giá bằng phần giá trị còn lại của MWG sau khi trừ BHX + ĐMX.'
], 14)

# Slide 4: valuation outputs
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Kết quả định giá base case')
add_table(slide, Inches(0.8), Inches(1.5), Inches(8.2), Inches(2.5), [
    ['Mảng','Giá trị (tỷ đồng)','% EV MWG'],
    ['BHX', f'{bhx_value_bn:,.1f}'.replace(',', '.'), f'{bhx_value_bn/mwg_ev_bn*100:.1f}%'],
    ['ĐMX', f'{dmx_value_bn:,.1f}'.replace(',', '.'), f'{dmx_value_bn/mwg_ev_bn*100:.1f}%'],
    ['TGDĐ', f'{tgdd_value_bn:,.1f}'.replace(',', '.'), f'{tgdd_value_bn/mwg_ev_bn*100:.1f}%'],
    ['Khác', f'{other_chains_bn:,.1f}'.replace(',', '.'), f'{other_chains_bn/mwg_ev_bn*100:.1f}%'],
    ['MWG EV', f'{mwg_ev_bn:,.1f}'.replace(',', '.'), '100.0%'],
])
add_bullets(slide, Inches(9.3), Inches(1.7), Inches(3.2), Inches(3.5), [
    f'BHX base: {bhx_value_bn:,.0f} tỷ'.replace(',', '.'),
    f'ĐMX base: {dmx_value_bn:,.0f} tỷ'.replace(',', '.'),
    f'TGDĐ base: {tgdd_value_bn:,.0f} tỷ'.replace(',', '.'),
    'ĐMX vẫn là core value lớn nhất trong base case.',
    'BHX đang đủ lớn để trở thành 1/3 EV của MWG nếu market tin vào profitability path.'
], 16)

# Slide 5: segment discussion
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Nhận định từng mảng')
add_bullets(slide, Inches(0.7), Inches(1.3), Inches(12), Inches(5.5), [
    'ĐMX: profit engine rõ nhất; có catalyst IPO; thêm optionality từ mảng dịch vụ thợ.',
    'BHX: đang chuyển từ story tăng trưởng doanh thu sang story lợi nhuận + leverage; định giá nhạy với biên ròng khoảng 3-5%.',
    'TGDĐ: mảng trưởng thành hơn, ít rerating hơn; vai trò chính là nền doanh thu/dòng tiền và phần giá trị còn lại trong SOTP.',
    'Nếu BHX duy trì biên khoảng 3% và mở rộng bền, upside valuation còn; nếu biên không giữ được thì BHX valuation dễ co lại.',
    'Base case hiện phù hợp hơn bull case vì bull case làm phần TGDĐ + other implied quá thấp nếu giữ nguyên EV MWG.'
], 18)

# Slide 6: caveats
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Điểm cần cẩn thận / Việc cần làm tiếp')
add_bullets(slide, Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.3), [
    'Chưa có standalone financial statements full-year sạch cho ĐMX và BHX trong phiên này.',
    'Net debt MWG và shares cần cập nhật từ BCTC / nguồn mới nhất nếu muốn target price chính xác hơn.',
    'TGDĐ đang là implied valuation, chưa phải tách riêng bằng mô hình độc lập hoàn chỉnh.',
    'Báo cáo này phù hợp để nhìn nhanh SOTP / tỷ trọng giá trị, chưa phải fairness value cuối cùng.',
    'Bước tiếp theo: cập nhật CĐKT/LCTT 2025, kéo thêm số standalone/IPO ĐMX và disclosure BHX để khóa valuation chắc hơn.'
], 18)

out = Path('MWG_BHX_DMX_SOTP_Report.pptx')
prs.save(out)
print(out.resolve())
