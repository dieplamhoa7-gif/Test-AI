import json, math, os, random
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

ROOT=Path('.')
DATA=ROOT/'data'
OUT=ROOT/'investment_strategy_rs_dashboard.pptx'

def load(name, default=None):
    try:
        return json.loads((DATA/name).read_text(encoding='utf-8'))
    except Exception:
        return default or {}

b4=load('b4_trend_pullback_locked_backtest.json')
bottom=load('v3_clean_split_rs_action_relaxed_a2_b2_backtest.json')
shake=load('shakeout_standard_deep_backtest.json')
sr=load('sr_quality_backtest_100.json')
current=load('strategy_results_cache.json')

prs=Presentation()
prs.slide_width=Inches(13.333)
prs.slide_height=Inches(7.5)
W,H=prs.slide_width, prs.slide_height

NAVY=RGBColor(8,18,39); NAVY2=RGBColor(12,30,65); CYAN=RGBColor(39,213,255); GREEN=RGBColor(58,214,146); RED=RGBColor(255,93,110); GOLD=RGBColor(255,196,87); WHITE=RGBColor(245,248,255); MUTED=RGBColor(164,178,205)

def bg(slide, title=None, subtitle=None):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb=NAVY
    # stock style grid
    for i in range(0,16):
        x=int(W*i/15); line=slide.shapes.add_connector(1,x,0,x,H); line.line.color.rgb=RGBColor(21,43,82); line.line.transparency=55; line.line.width=Pt(.5)
    for i in range(0,9):
        y=int(H*i/8); line=slide.shapes.add_connector(1,0,y,W,y); line.line.color.rgb=RGBColor(21,43,82); line.line.transparency=55; line.line.width=Pt(.5)
    # stylized price line
    pts=[]; base=float(H)*0.70
    for i in range(12):
        x=float(W)*(0.03+i*0.085); y=base + math.sin(i*1.1)*350 - i*22 + random.choice([-60,30,10])
        pts.append((int(x), int(max(H*.22,min(H*.88,y)))))
    for a,b in zip(pts,pts[1:]):
        line=slide.shapes.add_connector(1,a[0],a[1],b[0],b[1]); line.line.color.rgb=RGBColor(28,105,115); line.line.transparency=40; line.line.width=Pt(2)
    if title:
        t=slide.shapes.add_textbox(Inches(.45), Inches(.25), Inches(12.2), Inches(.55))
        p=t.text_frame.paragraphs[0]; p.text=title; p.font.bold=True; p.font.size=Pt(26); p.font.color.rgb=WHITE
    if subtitle:
        st=slide.shapes.add_textbox(Inches(.48), Inches(.82), Inches(12), Inches(.32))
        p=st.text_frame.paragraphs[0]; p.text=subtitle; p.font.size=Pt(10.5); p.font.color.rgb=MUTED

def add_box(slide,x,y,w,h,title,body='',color=CYAN):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(13,32,68); sh.fill.transparency=8
    sh.line.color.rgb=color; sh.line.transparency=25; sh.line.width=Pt(1.2)
    tf=sh.text_frame; tf.margin_left=Inches(.13); tf.margin_right=Inches(.1); tf.margin_top=Inches(.08)
    p=tf.paragraphs[0]; p.text=title; p.font.bold=True; p.font.size=Pt(15); p.font.color.rgb=color
    if body:
        for line in body.split('\n'):
            q=tf.add_paragraph(); q.text=line; q.font.size=Pt(10.5); q.font.color.rgb=WHITE; q.level=0
    return sh

def metric(slide,x,y,w,h,label,value,color=GREEN):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(15,37,77); sh.line.color.rgb=color; sh.line.transparency=25
    tf=sh.text_frame; tf.clear(); tf.margin_top=Inches(.08)
    p=tf.paragraphs[0]; p.text=str(value); p.alignment=PP_ALIGN.CENTER; p.font.bold=True; p.font.size=Pt(24); p.font.color.rgb=color
    q=tf.add_paragraph(); q.text=label; q.alignment=PP_ALIGN.CENTER; q.font.size=Pt(9.5); q.font.color.rgb=MUTED

def bullet(slide,x,y,w,h,items,font=13):
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.clear()
    for idx,it in enumerate(items):
        p=tf.paragraphs[0] if idx==0 else tf.add_paragraph(); p.text='• '+it; p.font.size=Pt(font); p.font.color.rgb=WHITE; p.space_after=Pt(5)
    return tb

def chart_bar(slide,x,y,w,h,cats,vals,title,color=GREEN):
    data=CategoryChartData(); data.categories=cats; data.add_series(title, vals)
    ch=slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    ch.has_legend=False; ch.chart_title.has_text_frame=True; ch.chart_title.text_frame.text=title
    ch.value_axis.tick_labels.font.size=Pt(8); ch.category_axis.tick_labels.font.size=Pt(8)
    try: ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb=color
    except Exception: pass
    return ch

def fmt(v):
    if v is None: return 'N/A'
    if isinstance(v,float): return f'{v:.2f}'
    return str(v)

# 1 cover
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.75), Inches(1.05), Inches(11.9), Inches(4.7)); box.fill.solid(); box.fill.fore_color.rgb=RGBColor(7,16,36); box.fill.transparency=7; box.line.color.rgb=CYAN
logo=DATA/'assets'/'lh-logo.jpg'
if logo.exists(): s.shapes.add_picture(str(logo), Inches(.95), Inches(1.25), width=Inches(1.1))
t=s.shapes.add_textbox(Inches(2.2),Inches(1.25),Inches(9.8),Inches(.8)); p=t.text_frame.paragraphs[0]; p.text='LH INVESTMENT'; p.font.size=Pt(36); p.font.bold=True; p.font.color.rgb=WHITE
st=s.shapes.add_textbox(Inches(2.22),Inches(2.0),Inches(9.6),Inches(1.2)); p=st.text_frame.paragraphs[0]; p.text='3 CHIẾN LƯỢC ĐẦU TƯ & MÔ HÌNH R/S'; p.font.size=Pt(30); p.font.bold=True; p.font.color.rgb=CYAN
bullet(s,2.25,3.08,8.7,1.2,['Tổng quan logic giao dịch','Định lượng chỉ báo kỹ thuật','Kết quả kiểm định / backtest','Khung xác định Support – Resistance'],14)

# 2 strategy overview
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'Tổng quan 3 chiến lược','Bộ lọc định lượng: R/S xác định vùng giá, chỉ báo xác nhận hành động giá')
add_box(s,.55,1.35,3.9,4.9,'1. Trend Pullback Pro','Mục tiêu: mua cổ phiếu đang có xu hướng tăng, hồi về hỗ trợ.\nBối cảnh: giá trên mây Ichimoku, RSI trung tính, MACD hồi.\nPhù hợp: trend còn khỏe, không mua đuổi.',CYAN)
add_box(s,4.72,1.35,3.9,4.9,'2. Support Rebound Hunter','Mục tiêu: bắt nhịp hồi tại hỗ trợ sau nhịp giảm/điều chỉnh.\nBối cảnh: gần hỗ trợ, RSI thấp, %BB thấp.\nĐiều kiện mới: WATCH cũng cần MACD histogram hồi.',GREEN)
add_box(s,8.88,1.35,3.9,4.9,'3. Shakeout Rebound','Mục tiêu: bắt cú rũ bỏ thủng hỗ trợ giả rồi reclaim.\nBối cảnh: breakdown 2–4.5%, sau đó hồi lại vùng hỗ trợ.\nPhù hợp: cổ phiếu có cú wash-out rõ.',GOLD)

# 3 indicators
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'Các nhóm chỉ báo sử dụng','Không dùng một chỉ báo đơn lẻ — chiến lược yêu cầu confluence')
add_box(s,.55,1.25,3.0,4.9,'Xu hướng','• Ichimoku cloud\n• MA20 / MA50 / MA200\n• Market structure\n• ROC20',CYAN)
add_box(s,3.8,1.25,3.0,4.9,'Động lượng','• RSI14\n• MACD line/signal\n• MACD histogram\n• ADX / DI',GREEN)
add_box(s,7.05,1.25,3.0,4.9,'Vị trí giá','• R/S zone\n• Distance to support\n• Bollinger %BB\n• Pivot / swing',GOLD)
add_box(s,10.3,1.25,2.55,4.9,'Rủi ro','• Stop theo support\n• Target theo resistance/RR\n• Bearish divergence filter\n• Volume confirmation',RED)

# 4 quantitative rules
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'Định lượng chỉ báo theo chiến lược','Các ngưỡng chính đang dùng trong pipeline local/output')
add_box(s,.55,1.2,4.0,5.3,'Trend Pullback Pro','Điều kiện lõi:\n• Giá trên Ichimoku cloud\n• Gần hỗ trợ: distance <= 3%\n• RSI: 48–62\n• MACD histogram hồi hoặc phân kỳ dương\n• Không có bearish divergence\n• Volume không quá xấu\nQuản trị:\n• Stop khoảng -6%\n• Chốt một phần tại target +6%',CYAN)
add_box(s,4.75,1.2,4.0,5.3,'Support Rebound Hunter','Hai nhánh A2/B2:\n• A2: gần hỗ trợ <= 2.5%\n• RSI <= 45\n• %BB <= 0.55\n• B2: above cloud + RSI 48–62\n• WATCH bắt buộc MACD hist recovering\n• Loại nếu bearish divergence',GREEN)
add_box(s,8.95,1.2,3.85,5.3,'Shakeout Rebound','Điều kiện lõi:\n• Thủng hỗ trợ cũ 2–4.5%\n• Chờ 1–3 nến reclaim/rebound\n• Entry nến sau xác nhận\n• Không mua nếu gap xấu\nQuản trị:\n• Stop động theo vùng hỗ trợ\n• Target ngắn 5–6% hoặc theo R/S',GOLD)

# 5 current signals
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'Tín hiệu hiện tại sau khi chạy VN100','Ngày chạy: 2026-05-02 — không có mã đạt BUY tuyệt đối')
metric(s,.7,1.25,2.2,1.1,'BUY toàn hệ thống',0,RED)
metric(s,3.15,1.25,2.2,1.1,'WATCH Trend Pullback',5,GREEN)
metric(s,5.6,1.25,2.2,1.1,'WATCH Bắt đáy',0,GOLD)
metric(s,8.05,1.25,2.2,1.1,'WATCH Shakeout',0,CYAN)
metric(s,10.5,1.25,2.2,1.1,'VN100 scanner',95,CYAN)
add_box(s,.75,2.8,11.8,3.5,'Danh sách WATCH đáng theo dõi','DXG — Mua 14.8 | Target 15.69 | SL 13.91 | Score 87.5\nFPT — Mua 73.4 | Target 77.8 | SL 69.0 | Score 75.0\nGEX — Mua 36.8 | Target 39.01 | SL 34.59 | Score 75.0\nMSN — Mua 74.3 | Target 78.76 | SL 69.84 | Score 75.0\nGAS — Mua 73.0 | Target 77.38 | SL 68.62 | Score 75.0',GREEN)

# 6 backtest strategy
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'Kết quả kiểm định 3 chiến lược','Số liệu lấy từ file backtest/output sẵn có trong data')
b4sum=b4.get('current180',{}).get('summary',{})
metric(s,.55,1.15,2.2,1.05,'Trend Pullback trades',b4sum.get('totalTrades','N/A'),CYAN)
metric(s,2.95,1.15,2.2,1.05,'Win rate',str(b4sum.get('winRatePct','N/A'))+'%',GREEN)
metric(s,5.35,1.15,2.2,1.05,'Avg PnL/trade',str(b4sum.get('avgPnlPct','N/A'))+'%',GOLD)
metric(s,7.75,1.15,2.2,1.05,'Sum PnL',str(b4sum.get('sumPnlPct','N/A'))+'%',GREEN)
metric(s,10.15,1.15,2.2,1.05,'Avg hold',str(b4sum.get('avgHoldSessions','N/A'))+' phiên',CYAN)
chart_bar(s,.65,2.65,4.0,3.7,['Wins','Losses'],[b4sum.get('wins',0), b4sum.get('losses',0)],'Trend Pullback: thắng/thua',GREEN)
# shakeout rows
rows=shake.get('rows',[])
cats=[]; vals=[]
for r in rows:
    name=r.get('group') or r.get('mode') or 'variant'
    oos=r.get('oos') or {}
    wr=oos.get('winRatePct') or oos.get('winRate') or 0
    cats.append(str(name)[:14]); vals.append(float(wr or 0))
if not cats: cats=['standard','deep']; vals=[0,0]
chart_bar(s,4.95,2.65,3.65,3.7,cats,vals,'Shakeout OOS win-rate (%)',GOLD)
# bottom window summary text
add_box(s,8.9,2.55,3.7,3.85,'Bắt đáy / Rebound','Backtest A2/B2 tách rõ:\n• R/S chỉ tạo vùng hỗ trợ/kháng cự, entry/stop/target\n• Action model dùng RSI, MACD, BB, volume, ROC, Ichimoku, divergence\n• Rule mới: loại WATCH nếu MACD histogram không hồi\n• Mục tiêu: giảm tín hiệu bắt dao rơi',RED)

# 7 RS method
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'Xác định Support / Resistance','R/S là mô hình confluence đa tín hiệu, không phải một pivot đơn')
add_box(s,.55,1.2,3.0,4.9,'Nguồn mức giá','• Swing high / swing low\n• Pivot high/low/close\n• MA20 / MA50 / MA200\n• VWAP\n• Donchian H/L/M\n• Fibonacci retracement',CYAN)
add_box(s,3.85,1.2,3.0,4.9,'Volume & vùng','• Volume-by-price clusters\n• Cụm thanh khoản\n• ATR để co giãn zone\n• Gộp các mức gần nhau\n• Chấm điểm strength',GREEN)
add_box(s,7.15,1.2,2.75,4.9,'Chọn active R/S','• Support dưới giá hiện tại\n• Resistance trên giá hiện tại\n• Ưu tiên gần + mạnh\n• Nếu thủng vùng thì chuyển support kế tiếp',GOLD)
add_box(s,10.15,1.2,2.65,4.9,'Ứng dụng','• Điểm mua\n• Vùng mua\n• Target\n• Stop-loss\n• RR filter\n• Không thay thế action signal',RED)

# 8 RS quantitative/backtest
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'Định lượng & kiểm định R/S','Kiểm tra chất lượng vùng hỗ trợ/kháng cự trên 100 mã')
ss=sr.get('summary',{})
side=ss.get('bySide',{})
sup=side.get('support',{}); res=side.get('resistance',{})
metric(s,.65,1.15,2.2,1.1,'Tổng event R/S',ss.get('totalEvents','N/A'),CYAN)
metric(s,3.05,1.15,2.2,1.1,'Support good-rate',str(sup.get('goodRatePct','N/A'))+'%',GREEN)
metric(s,5.45,1.15,2.2,1.1,'Support not-broken',str(sup.get('notBrokenRatePct','N/A'))+'%',GOLD)
metric(s,7.85,1.15,2.2,1.1,'Resistance good-rate',str(res.get('goodRatePct','N/A'))+'%',RED)
metric(s,10.25,1.15,2.2,1.1,'Universe',len(sr.get('symbols',[])) or 100,CYAN)
buckets=ss.get('byScoreBucket',{})
cats=list(buckets.keys()) or ['45-60','60-75','75-90','90+']
vals=[buckets.get(c,{}).get('goodRatePct',0) for c in cats]
chart_bar(s,.75,2.75,5.3,3.5,cats,vals,'Good-rate theo score bucket (%)',GREEN)
grades=ss.get('byGrade',{})
cats2=list(grades.keys()) or ['watch','valid','hard']; vals2=[grades.get(c,{}).get('notBrokenRatePct',0) for c in cats2]
chart_bar(s,6.65,2.75,5.3,3.5,cats2,vals2,'Not-broken rate theo grade (%)',CYAN)

# 9 visual flow
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'Luồng ra quyết định','Từ dữ liệu giá → R/S → chỉ báo → quyết định BUY/WATCH/REJECT')
steps=[('OHLCV','Giá, volume, lịch sử'),('R/S Engine','Cluster support/resistance'),('Indicator Layer','RSI, MACD, BB, ADX, MA, Ichimoku'),('Strategy Rules','Định lượng theo 3 chiến lược'),('Risk Plan','Entry, target, stop, RR'),('Output Cache','Web chỉ đọc JSON output')]
for i,(a,b) in enumerate(steps):
    x=.55+i*2.1
    add_box(s,x,2.1,1.75,2.25,a,b,[CYAN,GREEN,GOLD,RED,CYAN,GREEN][i])
    if i<5:
        arr=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+1.72), Inches(2.85), Inches(.45), Inches(.55)); arr.fill.solid(); arr.fill.fore_color.rgb=MUTED; arr.line.color.rgb=MUTED
bullet(s,1.0,5.25,11.6,1.1,['Nguyên tắc vận hành: mô hình và scan chạy local, xuất output JSON; web/Firebase/Render chỉ hiển thị cache, không tính toán nặng realtime.'],15)

# 10 conclusion
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'Kết luận & cách sử dụng','Chiến lược ưu tiên kỷ luật điểm mua, không mua đuổi')
add_box(s,.8,1.25,3.7,4.7,'Ưu tiên BUY','Chỉ mua khi đạt đủ bộ điều kiện: vùng giá hợp lý + động lượng xác nhận + rủi ro chấp nhận được.',GREEN)
add_box(s,4.85,1.25,3.7,4.7,'WATCH là danh sách canh','WATCH không phải lệnh mua. Dùng để theo dõi khi còn thiếu 1–2 điều kiện.',GOLD)
add_box(s,8.9,1.25,3.7,4.7,'R/S là khung giá','R/S giúp xác định entry, target, stop. Quyết định cuối cùng vẫn cần action model xác nhận.',CYAN)

# footer all slides
for idx,sl in enumerate(prs.slides,1):
    ft=sl.shapes.add_textbox(Inches(.45), Inches(7.12), Inches(12.4), Inches(.25))
    p=ft.text_frame.paragraphs[0]; p.text=f'LH INVESTMENT  |  Strategy & R/S Quant Framework  |  Slide {idx}/{len(prs.slides)}'; p.font.size=Pt(8.5); p.font.color.rgb=MUTED; p.alignment=PP_ALIGN.RIGHT

prs.save(OUT)
print(OUT)
